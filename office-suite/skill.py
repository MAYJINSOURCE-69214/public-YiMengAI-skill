#!/usr/bin/env python3
"""
忆梦 AI Office 五件套 v3.4.0
支持 Word (.docx, .doc)、Excel (.xlsx)、PowerPoint (.pptx)、PDF (.pdf)

调用方式：
    1. Python 导入：from skill import OfficeSuite; suite = OfficeSuite()
    2. 命令行：python skill.py read <文件路径>
"""
__version__ = "3.4.0"

import sys
import os
from pathlib import Path
from datetime import datetime
import importlib

# Word
from docx import Document as DocxDoc
from docx.shared import Inches, Pt, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

# Excel
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# PowerPoint
from pptx import Presentation as PptxPres
from pptx.util import Inches as PptxInches, Pt as PptxPt

# PDF
try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

# Word COM (for .doc)
try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False


class OfficeSuite:
    """Office 五件套操作类（支持 Python 导入调用）"""

    def __init__(self):
        self.supported_formats = {'.docx', '.doc', '.xlsx', '.pptx', '.pdf'}

    def _check_file(self, file_path):
        """检查文件是否存在"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在：{file_path}")
        ext = Path(file_path).suffix.lower()
        if ext not in self.supported_formats:
            raise ValueError(f"不支持的格式：{ext}，支持：{self.supported_formats}")
        return ext

    def read_docx(self, file_path):
        """读取 Word 文档，返回结构化数据"""
        ext = self._check_file(file_path)
        if ext != '.docx':
            raise ValueError(f"不是 Word 文件：{file_path}")

        doc = DocxDoc(file_path)
        result = {
            'file': file_path,
            'paragraphs': [],
            'tables': [],
            'comments': []
        }

        for para in doc.paragraphs:
            if para.text.strip():
                result['paragraphs'].append({
                    'style': para.style.name,
                    'text': para.text
                })

        for t_idx, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            result['tables'].append({
                'index': t_idx + 1,
                'rows': len(table.rows),
                'cols': len(table.columns),
                'data': table_data
            })

        for comment in doc.comments:
            result['comments'].append({
                'author': comment.author,
                'date': str(comment.date),
                'text': comment.text
            })

        return result

    def read_doc(self, file_path):
        """读取旧版 Word (.doc) 文档，返回结构化数据"""
        if not HAS_WIN32COM:
            raise ImportError("读取 .doc 文件需要 pywin32 库，请运行: pip install pywin32")

        # 转换为 .docx
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(file_path))
        
        # 临时文件
        temp_docx = file_path + "x" 
        doc.SaveAs2(temp_docx, FileFormat=16) # 16 = wdFormatXMLDocument
        doc.Close()
        word.Quit()
        
        # 读取转换后的 docx
        result = self.read_docx(temp_docx)
        
        # 清理临时文件
        if os.path.exists(temp_docx):
            os.remove(temp_docx)
            
        return result

    def read_xlsx(self, file_path, sheet_name=None, max_rows=500):
        """读取 Excel 文件，返回结构化数据（支持中文路径）"""
        ext = self._check_file(file_path)
        if ext != '.xlsx':
            raise ValueError(f"不是 Excel 文件：{file_path}")

        wb = load_workbook(file_path, data_only=True)
        result = {
            'file': file_path,
            'sheets': [],
            'sheet_names': wb.sheetnames
        }

        sheets_to_read = [sheet_name] if sheet_name else wb.sheetnames

        for name in sheets_to_read:
            if name not in wb.sheetnames:
                continue
            ws = wb[name]
            sheet_data = {
                'name': name,
                'rows': ws.max_row,
                'cols': ws.max_column,
                'data': []
            }

            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, max_rows), values_only=True):
                if any(cell is not None for cell in row):
                    sheet_data['data'].append(list(row))

            result['sheets'].append(sheet_data)

        return result

    def read_pptx(self, file_path):
        """读取 PowerPoint 文件，返回结构化数据"""
        ext = self._check_file(file_path)
        if ext != '.pptx':
            raise ValueError(f"不是 PowerPoint 文件：{file_path}")

        prs = PptxPres(file_path)
        result = {
            'file': file_path,
            'slides': [],
            'slide_count': len(prs.slides)
        }

        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                'index': i,
                'text': [],
                'tables': [],
                'images': []
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_data['text'].append(para.text)
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    slide_data['tables'].append(table_data)

            result['slides'].append(slide_data)

        return result

    def read_pdf(self, file_path):
        """读取 PDF 文件，返回结构化数据（文本型 PDF）"""
        if not HAS_PYPDF:
            raise ImportError("缺少 pypdf 库，请先安装：pip install pypdf")

        ext = self._check_file(file_path)
        if ext != '.pdf':
            raise ValueError(f"不是 PDF 文件：{file_path}")

        try:
            reader = PdfReader(file_path)
            if reader.is_encrypted:
                raise ValueError("该 PDF 文件已加密，请输入密码或联系发送方解除加密。")
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"PDF 读取失败：{e}")

        result = {
            'file': file_path,
            'pages': [],
            'page_count': len(reader.pages),
            'metadata': reader.metadata
        }

        for i, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                result['pages'].append({
                    'index': i,
                    'text': text or ''
                })
            except Exception:
                result['pages'].append({
                    'index': i,
                    'text': '（该页无法提取文本）'
                })

        return result

    def create_docx(self, output_path, title="文档标题", paragraphs=None):
        """创建 Word 文档"""
        doc = DocxDoc()
        doc.add_heading(title, level=1)
        doc.add_paragraph(f'创建时间：{datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')

        if paragraphs:
            for p in paragraphs:
                doc.add_paragraph(p)

        doc.save(output_path)
        return output_path

    def create_xlsx(self, output_path, sheets_data=None):
        """创建 Excel 文件"""
        wb = Workbook()

        if sheets_data:
            for i, (name, rows) in enumerate(sheets_data.items()):
                if i == 0:
                    ws = wb.active
                    ws.title = name
                else:
                    ws = wb.create_sheet(name)

                for row in rows:
                    ws.append(row)
        else:
            wb.active.title = "Sheet1"

        wb.save(output_path)
        return output_path

    def get_info(self, file_path):
        """获取文档信息摘要"""
        ext = self._check_file(file_path)

        if ext == '.docx':
            doc = DocxDoc(file_path)
            return {
                'type': 'Word',
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'comments': len(doc.comments),
                'non_empty_paragraphs': sum(1 for p in doc.paragraphs if p.text.strip())
            }
        elif ext == '.doc':
            # For .doc, we need to convert first or just return basic info
            # To keep it simple, we'll convert and read
            if HAS_WIN32COM:
                result = self.read_doc(file_path)
                return {
                    'type': 'Word (Legacy)',
                    'paragraphs': len(result['paragraphs']),
                    'tables': len(result['tables']),
                    'comments': len(result['comments'])
                }
            else:
                return {'type': 'Word (Legacy)', 'error': 'pywin32 not installed'}
        elif ext == '.xlsx':
            wb = load_workbook(file_path, data_only=True)
            sheets_info = {}
            for name in wb.sheetnames:
                ws = wb[name]
                sheets_info[name] = {'rows': ws.max_row, 'cols': ws.max_column}
            return {
                'type': 'Excel',
                'sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
                'sheets_detail': sheets_info
            }
        elif ext == '.pptx':
            prs = PptxPres(file_path)
            return {
                'type': 'PowerPoint',
                'slides': len(prs.slides)
            }
        elif ext == '.pdf':
            if not HAS_PYPDF:
                return {'type': 'PDF', 'error': '缺少 pypdf 库'}
            try:
                reader = PdfReader(file_path)
                return {
                    'type': 'PDF',
                    'pages': len(reader.pages),
                    'metadata': str(reader.metadata) if reader.metadata else '无'
                }
            except Exception as e:
                return {'type': 'PDF', 'error': str(e)}


# ============================================================
# 命令行兼容（原有功能保留）
# ============================================================

def read_word(file_path):
    suite = OfficeSuite()
    ext = Path(file_path).suffix.lower()
    if ext == '.doc':
        result = suite.read_doc(file_path)
    else:
        result = suite.read_docx(file_path)
        
    print(f"📘 Word 文档：{file_path}")
    print(f"📊 段落数：{len(result['paragraphs'])}")
    print(f"📋 表格数：{len(result['tables'])}")
    if result.get('comments'):
        print(f"💬 批注数：{len(result['comments'])}")
    print("=" * 80)
    for p in result['paragraphs']:
        print(f"[{p['style']}] {p['text'][:150]}")
    for t in result['tables']:
        print(f"\n📋 表格 {t['index']} ({t['rows']}行 x {t['cols']}列)")
        for row in t['data']:
            print(" | ".join([str(c)[:30] for c in row]))


def read_excel(file_path):
    suite = OfficeSuite()
    result = suite.read_xlsx(file_path)
    print(f" Excel 工作簿：{file_path}")
    print(f"📑 工作表数：{len(result['sheets'])}")
    print("=" * 80)
    for s in result['sheets']:
        print(f"\n📑 工作表：{s['name']} ({s['rows']}行 x {s['cols']}列)")
        for row in s['data'][:50]:
            print(" | ".join([str(c) if c else '' for c in row[:10]]))


def read_pptx(file_path):
    suite = OfficeSuite()
    result = suite.read_pptx(file_path)
    print(f"📙 PowerPoint：{file_path}")
    print(f"📊 幻灯片数：{result['slide_count']}")
    print("=" * 80)
    for s in result['slides']:
        print(f"\n 幻灯片 {s['index']}")
        for text in s['text']:
            print(f"  {text[:150]}")


def read_pdf(file_path):
    suite = OfficeSuite()
    result = suite.read_pdf(file_path)
    print(f"📕 PDF 文档：{file_path}")
    print(f"📄 页数：{result['page_count']}")
    if result['metadata']:
        print(f" 元数据：{result['metadata']}")
    print("=" * 80)
    for page in result['pages']:
        print(f"\n 第 {page['index']} 页")
        print(page['text'][:500] if page['text'] else '（无文本内容）')


def document_info(file_path):
    suite = OfficeSuite()
    info = suite.get_info(file_path)
    print(f"📄 类型：{info['type']}")
    for k, v in info.items():
        if k != 'type':
            print(f"  {k}: {v}")


def check_dependencies():
    """检查环境依赖"""
    print(f" 检查 Office 技能 v{__version__} 依赖...")
    packages = {
        'python-docx': 'docx',
        'openpyxl': 'openpyxl',
        'python-pptx': 'pptx',
        'pypdf': 'pypdf',
        'pywin32': 'win32com.client'
    }

    all_ok = True
    for name, mod in packages.items():
        try:
            importlib.import_module(mod)
            print(f"  ✅ {name}")
        except ImportError:
            print(f"  ❌ {name} (未安装)")
            all_ok = False

    if all_ok:
        print("\n✨ 所有依赖已就绪！")
    else:
        print("\n💡 安装命令：pip install python-docx openpyxl python-pptx pypdf pywin32")


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    command = sys.argv[1].lower()
    suite = OfficeSuite()

    if command == 'check':
        check_dependencies()
        sys.exit(0)

    if len(sys.argv) < 3:
        print(f"❌ 需要文件路径参数")
        sys.exit(1)

    file_path = sys.argv[2]

    if command == 'read':
        ext = Path(file_path).suffix.lower()
        if ext == '.docx' or ext == '.doc':
            read_word(file_path)
        elif ext == '.xlsx':
            read_excel(file_path)
        elif ext == '.pptx':
            read_pptx(file_path)
        elif ext == '.pdf':
            read_pdf(file_path)
        else:
            print(f"❌ 不支持的格式：{ext}")
    elif command == 'info':
        document_info(file_path)
    else:
        print(f" 未知命令：{command}")


if __name__ == '__main__':
    main()

# ========== 技能统计自动更新 ==========
def update_stats():
    import subprocess
    try:
        qwen_dir = Path(__file__).parent.parent.parent
        subprocess.run(
            [sys.executable, str(qwen_dir / "skills" / "update_stats.py"), "office-suite", __version__],
            capture_output=True, text=True, cwd=qwen_dir
        )
    except:
        pass

update_stats()
